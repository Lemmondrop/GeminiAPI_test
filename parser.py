import os
import pandas as pd
import traceback

# [PDF 파싱]
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

# [Word 파싱]
try:
    import docx
except ImportError:
    docx = None

# [PPT 파싱]
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

def extract_text_from_pdf(file_path):
    if not fitz: return "PyMuPDF 라이브러리가 설치되지 않았습니다."
    try:
        doc = fitz.open(file_path)
        content = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            content.append(f"\n--- [PDF Page {page_num + 1}] ---\n{page.get_text('text')}")
        doc.close()
        return "\n".join(content)
    except Exception as e:
        return f"[PDF 파싱 오류] {e}"

def extract_text_from_docx(file_path):
    if not docx: return "python-docx 라이브러리가 설치되지 않았습니다."
    try:
        doc = docx.Document(file_path)
        content = []
        # 단락 추출
        for para in doc.paragraphs:
            if para.text.strip():
                content.append(para.text.strip())
        # 표 데이터 추출
        for table in doc.tables:
            for row in table.rows:
                row_data = [cell.text.strip().replace('\n', ' ') for cell in row.cells if cell.text.strip()]
                if row_data:
                    content.append(" | ".join(row_data))
        return "\n".join(content)
    except Exception as e:
        return f"[DOCX 파싱 오류] {e}"

def extract_text_from_pptx(file_path):
    if not Presentation: return "python-pptx 라이브러리가 설치되지 않았습니다."
    try:
        prs = Presentation(file_path)
        content = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                content.append(f"\n--- [PPT Slide {i + 1}] ---\n" + "\n".join(slide_texts))
        return "\n".join(content)
    except Exception as e:
        return f"[PPTX 파싱 오류] {e}"

def extract_text_from_excel(file_path):
    try:
        # 모든 시트를 읽어서 Dictionary 형태로 가져옴
        dfs = pd.read_excel(file_path, sheet_name=None)
        content = []
        for sheet_name, df in dfs.items():
            content.append(f"\n--- [Excel Sheet: {sheet_name}] ---")
            # NaN 값 정리 및 Markdown 표 형태로 변환 (LLM이 읽기 매우 좋음)
            content.append(df.fillna("").to_markdown(index=False))
        return "\n".join(content)
    except Exception as e:
        return f"[Excel 파싱 오류] {e}"

def extract_text_from_csv(file_path):
    try:
        encodings = ['utf-8', 'utf-8-sig', 'cp949', 'euc-kr']
        for enc in encodings:
            try:
                df = pd.read_csv(file_path, encoding=enc)
                return f"\n--- [CSV Data] ---\n" + df.fillna("").to_markdown(index=False)
            except UnicodeDecodeError:
                continue
        return "[CSV 파싱 오류] 지원하지 않는 인코딩입니다."
    except Exception as e:
        return f"[CSV 파싱 오류] {e}"

def parse_any_file(file_path):
    """파일 확장자를 감지하여 알맞은 텍스트 추출기를 호출합니다."""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf': return extract_text_from_pdf(file_path)
    elif ext == '.docx': return extract_text_from_docx(file_path)
    elif ext == '.pptx': return extract_text_from_pptx(file_path)
    elif ext in ['.xlsx', '.xls']: return extract_text_from_excel(file_path)
    elif ext == '.csv': return extract_text_from_csv(file_path)
    else: return f"지원하지 않는 확장자입니다: {ext}"

if __name__ == "__main__":
    # 개별 테스트용
    print("모든 파서 로드 완료.")